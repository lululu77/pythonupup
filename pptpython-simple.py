# coding:utf-8
import sys
from pptx.dml.color import RGBColor
from pptx import Presentation
fr = open(str(sys.argv[1]),'r')
lines = fr.readlines()
SLD_LAYOUT_TITLE_AND_CONTENT = 0 
prs = Presentation()
title_layout = prs.slide_layouts[SLD_LAYOUT_TITLE_AND_CONTENT]
pptxname = str(sys.argv[1]).replace('txt', 'pptx')
for line in lines:
#	print("i is {}\n".format(i))
    if '  ' in line:
        newtext = line.split('  ')
        #print(newtext[1])
        slide =prs.slides.add_slide(title_layout)
        title = slide.shapes.title
        title.text = newtext[0]
        subtitle = slide.placeholders[1]
        tf = subtitle.text_frame
        nomaltext = newtext[1]
        p = tf.paragraphs[0]
        p.text = ''
        run = p.add_run()
        run.text = nomaltext
        run.font.color.rgb = RGBColor(0xFF, 0x00, 0x00)
	
prs.save(pptxname)

