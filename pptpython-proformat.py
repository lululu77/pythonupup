# coding:utf-8
import sys
from pptx import Presentation
fr = open(str(sys.argv[1]),'r')
lines = fr.readlines()
SLD_LAYOUT_TITLE_AND_CONTENT = 0 
prs = Presentation()
title_layout = prs.slide_layouts[SLD_LAYOUT_TITLE_AND_CONTENT]

for line in lines:
#	print("i is {}\n".format(i))
	if '(' in line:
		newtext = line.split('(')
	slide =prs.slides.add_slide(title_layout)
	title = slide.shapes.title
	title.text = newtext[0]
	subtitle = slide.placeholders[1]
	subtitle.text = '('+newtext[1]
prs.save('test.pptx')

